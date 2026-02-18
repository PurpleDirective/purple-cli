"""
Purple Docs Server
MCP server for reading and creating documents.
Reads: PDF, Excel. Creates: Excel (.xlsx), Word (.docx), PowerPoint (.pptx).
"""

import json
import math
from datetime import datetime
from pathlib import Path
from fastmcp import FastMCP

mcp = FastMCP("purple-docs")

# Restrict all file operations to these directory trees
PERMITTED_ROOTS = [
    Path.home().resolve(),
    Path("/tmp").resolve(),
]

# Auto-route documents into organized subfolders
SUBFOLDERS = {
    "excel": "excel",
    "word": "word",
    "powerpoint": "powerpoint",
}


def _sanitize_filename(filename: str) -> str:
    """Strip path separators from filenames to prevent directory traversal."""
    return Path(filename).name


def _validate_path(path: Path, operation: str) -> str | None:
    """Check that a resolved path falls under PERMITTED_ROOTS.
    Returns None if valid, or an error string if the path is outside allowed roots."""
    resolved = path.resolve()
    for root in PERMITTED_ROOTS:
        if resolved == root or str(resolved).startswith(str(root) + "/"):
            return None
    return f"Error: {operation} blocked -- path '{resolved}' is outside permitted directories ({', '.join(str(r) for r in PERMITTED_ROOTS)})"


def _resolve_output(output_dir: str, doc_type: str) -> Path | str:
    """Route files into type-specific subfolders. Creates directory if needed.
    Skips adding subfolder if it already appears in the path.
    Returns an error string if the path is outside permitted roots."""
    subfolder = SUBFOLDERS[doc_type]
    base = Path(output_dir).expanduser()
    if subfolder not in [p.lower() for p in base.parts]:
        base = base / subfolder
    err = _validate_path(base, f"create_{doc_type}")
    if err:
        return err
    base.mkdir(parents=True, exist_ok=True)
    return base


@mcp.tool()
def list_directory(
    directory_path: str,
    pattern: str = "*"
) -> str:
    """List files in a directory. Returns filenames, sizes, and types. Use this BEFORE reading files to see what exists. pattern can be '*' for all files, '*.pdf' for PDFs only, etc."""
    path = Path(directory_path).expanduser()
    err = _validate_path(path, "list_directory")
    if err:
        return err
    if not path.exists():
        return f"Error: directory not found: {path}"
    if not path.is_dir():
        return f"Error: not a directory: {path}"

    entries = []
    try:
        for item in sorted(path.glob(pattern)):
            if item.name.startswith("."):
                continue
            if item.is_dir():
                count = sum(1 for _ in item.iterdir() if not _.name.startswith("."))
                entries.append(f"  [DIR]  {item.name}/ ({count} items)")
            else:
                size = item.stat().st_size
                if size < 1024:
                    size_str = f"{size} B"
                elif size < 1024 * 1024:
                    size_str = f"{size / 1024:.1f} KB"
                else:
                    size_str = f"{size / (1024 * 1024):.1f} MB"
                entries.append(f"  {size_str:>10}  {item.name}")
    except PermissionError:
        return f"Error: permission denied: {path}"

    if not entries:
        return f"Directory {path} is empty (pattern: {pattern})"

    total = len(entries)
    if total > 100:
        entries = entries[:100]
        entries.append(f"\n... truncated ({total} total items). Use pattern parameter to filter.")

    return f"Directory: {path} ({total} items, pattern: {pattern})\n" + "\n".join(entries)


@mcp.tool()
def read_excel(
    file_path: str,
    sheet_name: str = "all"
) -> str:
    """Read an Excel file and return its contents as text. sheet_name can be 'all' or a specific sheet name."""
    from openpyxl import load_workbook

    path = Path(file_path).expanduser()
    err = _validate_path(path, "read_excel")
    if err:
        return err
    if not path.exists():
        return f"Error: file not found: {path}"

    try:
        wb = load_workbook(str(path), data_only=True)
    except Exception as e:
        return f"Error opening Excel file: {e}"

    sheets = wb.sheetnames if sheet_name == "all" else [sheet_name]
    output = [f"Excel: {path.name} ({len(wb.sheetnames)} sheets: {', '.join(wb.sheetnames)})"]

    for name in sheets:
        if name not in wb.sheetnames:
            output.append(f"\nSheet '{name}' not found")
            continue
        ws = wb[name]
        output.append(f"\n=== {name} ({ws.max_row} rows x {ws.max_column} cols) ===")
        for row in ws.iter_rows(values_only=True):
            output.append(" | ".join(str(cell if cell is not None else "") for cell in row))

    wb.close()

    result = "\n".join(output)
    if len(result) > 30000:
        result = result[:30000] + "\n\n... truncated. Use sheet_name parameter to read specific sheets."
    return result


@mcp.tool()
def read_pdf(
    file_path: str,
    pages: str = "all"
) -> str:
    """Read a PDF file and extract text and tables. pages can be 'all', a single number like '1', or a range like '1-5'. Returns extracted text and any tables found."""
    import pdfplumber

    path = Path(file_path).expanduser()
    err = _validate_path(path, "read_pdf")
    if err:
        return err
    if not path.exists():
        return f"Error: file not found: {path}"
    if not path.suffix.lower() == ".pdf":
        return f"Error: not a PDF file: {path.name}"

    try:
        pdf = pdfplumber.open(str(path))
    except Exception as e:
        return f"Error opening PDF: {e}"

    total_pages = len(pdf.pages)

    try:
        if pages == "all":
            page_nums = list(range(total_pages))
        elif "-" in pages:
            start, end = pages.split("-", 1)
            start = max(0, int(start) - 1)
            end = min(total_pages, int(end))
            page_nums = list(range(start, end))
        else:
            page_nums = [int(pages) - 1]
    except (ValueError, TypeError):
        pdf.close()
        return f"Error: invalid pages parameter '{pages}'. Use 'all', a number like '3', or a range like '1-5'."

    page_nums = [n for n in page_nums if 0 <= n < total_pages]
    if not page_nums:
        pdf.close()
        return f"Error: no valid pages in range. PDF has {total_pages} pages (use 1-{total_pages})."

    output = [f"PDF: {path.name} ({total_pages} pages)"]

    for i in page_nums:
        page = pdf.pages[i]
        output.append(f"\n--- Page {i + 1} ---")

        text = page.extract_text()
        if text:
            output.append(text)

        tables = page.extract_tables()
        for t_idx, table in enumerate(tables):
            output.append(f"\n[Table {t_idx + 1}]")
            for row in table:
                output.append(" | ".join(str(cell or "") for cell in row))

    pdf.close()

    result = "\n".join(output)
    if len(result) > 30000:
        result = result[:30000] + f"\n\n... truncated ({len(result)} chars total). Use pages parameter to read specific pages."
    return result


@mcp.tool()
def create_excel(
    filename: str,
    sheets: str,
    output_dir: str = "."
) -> str:
    """Create an Excel workbook. sheets is a JSON string: [{"name": "Sheet1", "headers": ["A","B"], "rows": [["val1","val2"]]}]"""
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    from openpyxl.utils import get_column_letter

    try:
        sheet_data = json.loads(sheets)
    except json.JSONDecodeError:
        return "Error: sheets must be valid JSON"

    if not isinstance(sheet_data, list):
        return "Error: sheets must be a JSON array"

    wb = Workbook()
    wb.remove(wb.active)

    header_fill = PatternFill(start_color="4472C4", end_color="4472C4", fill_type="solid")
    header_font_white = Font(bold=True, size=12, color="FFFFFF")
    thin_border = Border(
        left=Side(style="thin"),
        right=Side(style="thin"),
        top=Side(style="thin"),
        bottom=Side(style="thin"),
    )

    for sheet in sheet_data:
        ws = wb.create_sheet(title=sheet.get("name", "Sheet"))

        headers = sheet.get("headers", [])
        if headers:
            for col, header in enumerate(headers, 1):
                cell = ws.cell(row=1, column=col, value=header)
                cell.font = header_font_white
                cell.fill = header_fill
                cell.alignment = Alignment(horizontal="center")
                cell.border = thin_border

        rows = sheet.get("rows", [])
        for row_idx, row in enumerate(rows, 2):
            for col_idx, value in enumerate(row, 1):
                cell = ws.cell(row=row_idx, column=col_idx, value=_convert_value(value))
                cell.border = thin_border

        for col in range(1, len(headers) + 1):
            max_len = max(
                len(str(ws.cell(row=r, column=col).value or ""))
                for r in range(1, len(rows) + 2)
            )
            ws.column_dimensions[get_column_letter(col)].width = min(max_len + 4, 50)

    filename = _sanitize_filename(filename)
    if not filename.endswith(".xlsx"):
        filename += ".xlsx"

    outdir = _resolve_output(output_dir, "excel")
    if isinstance(outdir, str):
        return outdir
    filepath = outdir / filename
    wb.save(str(filepath))
    return f"Created Excel: {filepath.resolve()} ({len(sheet_data)} sheets, {sum(len(s.get('rows', [])) for s in sheet_data)} rows)"


@mcp.tool()
def create_word(
    filename: str,
    title: str,
    content: str,
    output_dir: str = "."
) -> str:
    """Create a Word document. content is a JSON string: [{"type": "heading", "text": "...", "level": 1}, {"type": "paragraph", "text": "..."}, {"type": "list", "items": ["a","b"]}, {"type": "table", "headers": ["A","B"], "rows": [["v1","v2"]]}]"""
    from docx import Document
    from docx.shared import Inches, Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    try:
        blocks = json.loads(content)
    except json.JSONDecodeError:
        return "Error: content must be valid JSON"

    if not isinstance(blocks, list):
        return "Error: content must be a JSON array"

    doc = Document()

    title_style = doc.styles["Title"]
    title_style.font.color.rgb = RGBColor(0x44, 0x72, 0xC4)
    doc.add_heading(title, level=0)

    for block in blocks:
        block_type = block.get("type", "paragraph")

        if block_type == "heading":
            level = block.get("level", 1)
            doc.add_heading(block.get("text", ""), level=level)

        elif block_type == "paragraph":
            p = doc.add_paragraph(block.get("text", ""))
            if block.get("bold"):
                for run in p.runs:
                    run.bold = True
            if block.get("italic"):
                for run in p.runs:
                    run.italic = True

        elif block_type == "list":
            items = block.get("items", [])
            list_style = block.get("style", "bullet")
            for item in items:
                if list_style == "numbered":
                    doc.add_paragraph(item, style="List Number")
                else:
                    doc.add_paragraph(item, style="List Bullet")

        elif block_type == "table":
            headers = block.get("headers", [])
            rows = block.get("rows", [])
            table = doc.add_table(rows=1 + len(rows), cols=len(headers))
            table.style = "Light Grid Accent 1"

            for i, header in enumerate(headers):
                table.rows[0].cells[i].text = str(header)

            for row_idx, row in enumerate(rows):
                for col_idx, value in enumerate(row[:len(headers)]):
                    table.rows[row_idx + 1].cells[col_idx].text = str(value)

        elif block_type == "page_break":
            doc.add_page_break()

    filename = _sanitize_filename(filename)
    if not filename.endswith(".docx"):
        filename += ".docx"

    outdir = _resolve_output(output_dir, "word")
    if isinstance(outdir, str):
        return outdir
    filepath = outdir / filename
    doc.save(str(filepath))
    return f"Created Word doc: {filepath.resolve()} ({len(blocks)} content blocks)"


@mcp.tool()
def create_powerpoint(
    filename: str,
    title: str,
    slides: str,
    output_dir: str = "."
) -> str:
    """Create a PowerPoint presentation. slides is a JSON string: [{"title": "Slide Title", "content": "Bullet text\\nMore text", "layout": "bullet"}]. Layouts: title, bullet, blank."""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    try:
        slide_data = json.loads(slides)
    except json.JSONDecodeError:
        return "Error: slides must be valid JSON"

    if not isinstance(slide_data, list):
        return "Error: slides must be a JSON array"

    prs = Presentation()

    title_layout = prs.slide_layouts[0]
    bullet_layout = prs.slide_layouts[1]
    blank_layout = prs.slide_layouts[6]

    slide = prs.slides.add_slide(title_layout)
    slide.shapes.title.text = title
    if 1 in slide.placeholders:
        slide.placeholders[1].text = datetime.now().strftime("%B %d, %Y")

    for s in slide_data:
        layout = s.get("layout", "bullet")

        if layout == "title":
            sl = prs.slides.add_slide(title_layout)
            sl.shapes.title.text = s.get("title", "")
            if len(sl.placeholders) > 1:
                sl.placeholders[1].text = s.get("content", "")

        elif layout == "bullet":
            sl = prs.slides.add_slide(bullet_layout)
            sl.shapes.title.text = s.get("title", "")
            body = sl.placeholders[1]
            tf = body.text_frame
            tf.text = ""
            lines = s.get("content", "").split("\n")
            for i, line in enumerate(lines):
                if i == 0:
                    tf.text = line.strip()
                else:
                    p = tf.add_paragraph()
                    p.text = line.strip()
                    p.level = 0

        elif layout == "blank":
            sl = prs.slides.add_slide(blank_layout)
            if s.get("title"):
                txBox = sl.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
                txBox.text_frame.text = s.get("title", "")

    filename = _sanitize_filename(filename)
    if not filename.endswith(".pptx"):
        filename += ".pptx"

    outdir = _resolve_output(output_dir, "powerpoint")
    if isinstance(outdir, str):
        return outdir
    filepath = outdir / filename
    prs.save(str(filepath))
    return f"Created PowerPoint: {filepath.resolve()} ({len(slide_data) + 1} slides)"


def _convert_value(value):
    """Convert string values to appropriate types for Excel.
    Preserves leading zeros (zip codes, IDs) as strings.
    Rejects NaN/inf. Handles lists/dicts gracefully."""
    if isinstance(value, (list, dict)):
        return str(value)
    if isinstance(value, (int, float)):
        if math.isnan(value) or math.isinf(value):
            return str(value)
        return value
    if isinstance(value, str):
        if len(value) > 1 and value[0] == "0" and value != "0" and "." not in value:
            return value
        try:
            return int(value)
        except ValueError:
            try:
                f = float(value)
                if math.isnan(f) or math.isinf(f):
                    return value
                return f
            except ValueError:
                return value
    return value


if __name__ == "__main__":
    mcp.run()
